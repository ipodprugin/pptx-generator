from pptx import Presentation  


async def render_text(input_path, model, output_path, jinja2_env):
    ppt = Presentation(input_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                await _render_text_frame(shape.text_frame, model, jinja2_env)
    ppt.save(output_path)


async def _render_text_frame(text_frame, model, jinja2_env):
    last_ok = True
    tmp_str = ""
    tmp_runs = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            cur_text = run.text
            tmp_str += cur_text
            tmp_runs.append(run)
            try:
                rtemplate = jinja2_env.from_string(tmp_str)
                rendered_text = rtemplate.render(model)

                for run in tmp_runs:
                    run.text = rendered_text
                    rendered_text = "" # overwrites text

                tmp_str = ""
                tmp_runs = []
                last_ok = True
            except Exception as e:
                #could not finish, i.e. have to append!
                last_ok = False

async def replace_images_by_shape_text(images: dict, template_path: str, output_path: str):
    prs = Presentation(template_path)
    for image in images:
        image_file = images[image]
        search_str = image
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text.find(search_str) != -1:
                        horiz_ = shape.left
                        vert_ = shape.top
                        height_ = shape.height
                        width_ = shape.width
                        slide.shapes.add_picture(image_file, horiz_, vert_, width_, height_)
    prs.save(output_path)
